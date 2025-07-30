"""
Утилиты для обработки файлов на основе agno playground.
Максимально нативная интеграция с agno.
"""

import pandas as pd
import io
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from typing import List, Optional
from fastapi import HTTPException, UploadFile
from agno.media import Audio, Image, Video, File as FileMedia
from agno.utils.log import logger


def determine_content_type_by_filename(filename: str) -> Optional[str]:
    """Определяет content-type по расширению файла"""
    if not filename:
        return None
    
    extension = filename.lower().split('.')[-1] if '.' in filename else None
    if not extension:
        return None
    
    # Маппинг расширений на MIME типы
    extension_map = {
        # Изображения
        'png': 'image/png',
        'jpg': 'image/jpeg', 
        'jpeg': 'image/jpeg',
        'webp': 'image/webp',
        
        # Аудио
        'wav': 'audio/wav',
        'mp3': 'audio/mp3',
        
        # Видео
        'mp4': 'video/mp4',
        'webm': 'video/webm',
        'mov': 'video/quicktime',
        
        # Документы
        'pdf': 'application/pdf',
        'csv': 'text/csv',
        'json': 'application/json',
        'txt': 'text/plain',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'xls': 'application/vnd.ms-excel',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    }
    
    return extension_map.get(extension)


def process_image(file: UploadFile) -> Image:
    """Обработка изображений как в agno playground"""
    content = file.file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty file")
    return Image(content=content)


def process_audio(file: UploadFile) -> Audio:
    """Обработка аудио как в agno playground"""
    content = file.file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty file")
    format = None
    if file.filename and "." in file.filename:
        format = file.filename.split(".")[-1].lower()
    elif file.content_type:
        format = file.content_type.split("/")[-1]

    return Audio(content=content, format=format)


def process_video(file: UploadFile) -> Video:
    """Обработка видео как в agno playground"""
    content = file.file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty file")
    return Video(content=content, format=file.content_type)


def process_document(file: UploadFile) -> Optional[FileMedia]:
    """Обработка документов как в agno playground"""
    try:
        content = file.file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty file")

        return FileMedia(content=content)
    except Exception as e:
        logger.error(f"Error processing document {file.filename}: {e}")
        return None


def process_csv(file: UploadFile) -> Optional[FileMedia]:
    """
    Специальная обработка CSV файлов для обхода ограничения OpenAI API.
    Конвертирует CSV в текстовое представление, удобное для LLM.
    """
    try:
        content = file.file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty CSV file")

        # Декодируем как текст
        csv_text = content.decode('utf-8')
        lines = csv_text.strip().split('\n')
        
        if not lines:
            raise HTTPException(status_code=400, detail="Empty CSV content")

        # Создаем читаемое представление для LLM
        headers = lines[0] if lines else ""
        data_rows = lines[1:] if len(lines) > 1 else []
        
        formatted_content = f"""CSV Table Data from file: {file.filename}

Headers: {headers}
Total rows: {len(data_rows)}
Data preview (first 10 rows):

{chr(10).join(data_rows[:10])}"""

        if len(data_rows) > 10:
            formatted_content += f"\n\n... and {len(data_rows) - 10} more rows"

        # Создаем FileMedia с mime_type text/plain чтобы OpenAI принял его
        return FileMedia(content=formatted_content.encode('utf-8'), mime_type='text/plain')
    except UnicodeDecodeError:
        logger.error(f"Failed to decode CSV file {file.filename} as UTF-8")
        return None
    except Exception as e:
        logger.error(f"Error processing CSV {file.filename}: {e}")
        return None


def convert_csv_to_text(file: UploadFile) -> str:
    """Конвертирует CSV файл в структурированный текст"""
    try:
        content = file.file.read()
        file.file.seek(0)  # Возвращаем указатель в начало
        df = pd.read_csv(io.StringIO(content.decode('utf-8')))
        
        text_content = f"""CSV Table: {file.filename}

Columns: {', '.join(df.columns.tolist())}
Rows: {len(df)}

Data Summary:
{df.describe(include='all').to_string() if not df.empty else 'No data'}

Full Data:
{df.to_string(index=False)}"""
        
        return text_content
    except Exception as e:
        logger.error(f"Error converting CSV {file.filename}: {e}")
        return f"CSV file: {file.filename} (conversion failed: {str(e)})"


def convert_excel_to_text(file: UploadFile) -> str:
    """Конвертирует Excel файл в структурированный текст"""
    try:
        content = file.file.read()
        file.file.seek(0)
        
        # Читаем все листы
        excel_data = pd.read_excel(io.BytesIO(content), sheet_name=None)
        
        text_content = f"Excel File: {file.filename}\n"
        text_content += f"Sheets: {len(excel_data)}\n\n"
        
        for sheet_name, df in excel_data.items():
            text_content += f"=== Sheet: {sheet_name} ===\n"
            text_content += f"Columns: {', '.join(df.columns.tolist())}\n"
            text_content += f"Rows: {len(df)}\n\n"
            
            if not df.empty:
                text_content += f"Data:\n{df.to_string(index=False)}\n\n"
            else:
                text_content += "No data in this sheet\n\n"
        
        return text_content
    except Exception as e:
        logger.error(f"Error converting Excel {file.filename}: {e}")
        return f"Excel file: {file.filename} (conversion failed: {str(e)})"


def convert_docx_to_text(file: UploadFile) -> str:
    """Конвертирует Word документ в текст"""
    try:
        content = file.file.read()
        file.file.seek(0)
        
        doc = Document(io.BytesIO(content))
        text_content = f"Word Document: {file.filename}\n\n"
        
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                text_content += f"{paragraph.text}\n"
        
        # Добавляем таблицы если есть
        if doc.tables:
            text_content += "\n=== Tables ===\n"
            for i, table in enumerate(doc.tables):
                text_content += f"\nTable {i+1}:\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    text_content += f"{row_text}\n"
        
        return text_content
    except Exception as e:
        logger.error(f"Error converting DOCX {file.filename}: {e}")
        return f"Word document: {file.filename} (conversion failed: {str(e)})"


def convert_pptx_to_text(file: UploadFile) -> str:
    """Конвертирует PowerPoint презентацию в текст"""
    try:
        content = file.file.read()
        file.file.seek(0)
        
        prs = Presentation(io.BytesIO(content))
        text_content = f"PowerPoint Presentation: {file.filename}\n"
        text_content += f"Slides: {len(prs.slides)}\n\n"
        
        for i, slide in enumerate(prs.slides):
            text_content += f"=== Slide {i+1} ===\n"
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_content += f"{shape.text}\n"
            
            text_content += "\n"
        
        return text_content
    except Exception as e:
        logger.error(f"Error converting PPTX {file.filename}: {e}")
        return f"PowerPoint file: {file.filename} (conversion failed: {str(e)})"


def convert_pdf_to_text(file: UploadFile) -> str:
    """Конвертирует PDF файл в текст"""
    try:
        content = file.file.read()
        file.file.seek(0)
        
        reader = PdfReader(io.BytesIO(content))
        text_content = f"PDF Document: {file.filename}\n"
        text_content += f"Pages: {len(reader.pages)}\n\n"
        
        for i, page in enumerate(reader.pages):
            text_content += f"=== Page {i+1} ===\n"
            try:
                text_content += page.extract_text()
                text_content += "\n\n"
            except Exception as e:
                text_content += f"(Page extraction failed: {str(e)})\n\n"
        
        return text_content
    except Exception as e:
        logger.error(f"Error converting PDF {file.filename}: {e}")
        return f"PDF file: {file.filename} (conversion failed: {str(e)})"


def process_files(files: Optional[List[UploadFile]]) -> tuple[List[Image], List[Audio], List[Video], List[FileMedia]]:
    """
    Обработка файлов как в agno playground.
    Возвращает кортеж (images, audios, videos, documents)
    """
    base64_images: List[Image] = []
    base64_audios: List[Audio] = []
    base64_videos: List[Video] = []
    input_files: List[FileMedia] = []

    if not files:
        return base64_images, base64_audios, base64_videos, input_files

    for file in files:
        # Определяем content_type: используем оригинальный или определяем по расширению
        content_type = file.content_type
        
        # Если content_type отсутствует или является generic типом, определяем по расширению
        if not content_type or content_type in ["application/octet-stream", "text/plain"]:
            filename_content_type = determine_content_type_by_filename(file.filename)
            if filename_content_type:
                content_type = filename_content_type
        
        # Обработка файла по типу
        
        # Изображения
        if content_type in ["image/png", "image/jpeg", "image/jpg", "image/webp"]:
            try:
                base64_image = process_image(file)
                base64_images.append(base64_image)
                logger.info(f"✅ Image processed: {file.filename}")
            except Exception as e:
                logger.error(f"Error processing image {file.filename}: {e}")
            continue
                
        # Аудио
        elif content_type in ["audio/wav", "audio/mp3", "audio/mpeg"]:
            try:
                base64_audio = process_audio(file)
                base64_audios.append(base64_audio)
                logger.info(f"✅ Audio processed: {file.filename}")
            except Exception as e:
                logger.error(f"Error processing audio {file.filename}: {e}")
            continue
                
        # Видео
        elif content_type in [
            "video/x-flv", "video/quicktime", "video/mpeg", "video/mpegs",
            "video/mpgs", "video/mpg", "video/mp4", "video/webm", "video/wmv", "video/3gpp"
        ]:
            try:
                base64_video = process_video(file)
                base64_videos.append(base64_video)
                logger.info(f"✅ Video processed: {file.filename}")
            except Exception as e:
                logger.error(f"Error processing video {file.filename}: {e}")
            continue
                
        # CSV файлы 
        elif content_type == "text/csv":
            try:
                text_content = convert_csv_to_text(file)
                text_file = FileMedia(content=text_content.encode('utf-8'))
                input_files.append(text_file)
                logger.info(f"✅ CSV converted to text: {file.filename}")
            except Exception as e:
                logger.error(f"Error converting CSV {file.filename}: {e}")
                continue
                
        # Excel файлы (.xlsx, .xls)
        elif content_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
            try:
                text_content = convert_excel_to_text(file)
                text_file = FileMedia(content=text_content.encode('utf-8'))
                input_files.append(text_file)
                logger.info(f"✅ Excel converted to text: {file.filename}")
            except Exception as e:
                logger.error(f"Error converting Excel {file.filename}: {e}")
                continue
                
        # Word документы (.docx)
        elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            try:
                text_content = convert_docx_to_text(file)
                text_file = FileMedia(content=text_content.encode('utf-8'))
                input_files.append(text_file)
                logger.info(f"✅ Word document converted to text: {file.filename}")
            except Exception as e:
                logger.error(f"Error converting Word document {file.filename}: {e}")
                continue
                
        # PowerPoint презентации (.pptx)
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            try:
                text_content = convert_pptx_to_text(file)
                text_file = FileMedia(content=text_content.encode('utf-8'))
                input_files.append(text_file)
                logger.info(f"✅ PowerPoint converted to text: {file.filename}")
            except Exception as e:
                logger.error(f"Error converting PowerPoint {file.filename}: {e}")
                continue
                
        # Остальные документы (PDF остается в нативной agno обработке)
        elif content_type in [
            "application/pdf", 
            "text/plain", "application/json"
        ]:
            document_file = process_document(file)
            if document_file is not None:
                input_files.append(document_file)
        else:
            # Если не удалось определить тип файла
            if content_type is None:
                logger.warning(f"Could not determine content-type for file: {file.filename}")
                raise HTTPException(status_code=400, detail=f"Could not determine file type for: {file.filename}")
            else:
                logger.warning(f"Unsupported content-type: {content_type} for file: {file.filename}")
                raise HTTPException(status_code=400, detail=f"Unsupported file type: {content_type}")

    return base64_images, base64_audios, base64_videos, input_files 